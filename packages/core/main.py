from fastapi import FastAPI
from sentence_transformers import SentenceTransformer
from PyPDF2 import PdfReader
from pptx import Presentation
import pytesseract
from PIL import Image

app = FastAPI()


@app.get("/")
async def root():
    return {"message": "Hello World"}


class PDFProcessor:
    def extract_text(self, file_path):
        text = ""
        reader = PdfReader(file_path)
        for i in range(len(reader.pages)):
            page = reader.pages[i]
            text += page.extract_text()
        
        return text
              


class PPTProcessor:
    def extract_text(self, file_path):
        text = ""
        pptx_reader = Presentation(file_path)
        for slide in pptx_reader.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text

        return text            
        

class ImageProcessor:
    def extract_text(self, file_path):
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)


class DocumentVectorizer:
    def __init__(self):
        self.model = SentenceTransformer('all-MiniLM-L6-v2')
        self.file_processors = {
            'pdf': PDFProcessor(),
            'ppt': PPTProcessor(),
            'jpg': ImageProcessor(),
        }

obj = ImageProcessor()
print(obj.extract_text("test.jpg"))